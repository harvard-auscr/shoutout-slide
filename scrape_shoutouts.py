"""Scrape the shout-out slide(s) out of every General Meeting deck in ``dataset/``.

Why this exists: the shout-out generator needs a corpus of real shout-out slides
(geometry, fonts, animation XML) to learn from and to use as a template. Each GM
deck is 30-90 slides; we only want the one or two that carry the shout-outs.

How a shout-out block is recognised (measured across all 24 decks):
    * a ``SECTION_HEADER`` slide whose only text is "Shoutouts", followed by
    * one or more consecutive ``ONE_COLUMN_TEXT`` slides full of textboxes.
The header itself is NOT kept -- the user wants "the shout-out slide itself".

How extraction works: instead of copying shapes into a fresh file (which loses
the layout, theme and, crucially, the per-shape click animations), we open the
deck and *delete every other slide*. python-pptx only serialises parts that are
still reachable from the presentation part, so the orphaned slides, their notes
and media vanish on save, while the kept slide carries its layout, master, theme
and ``<p:timing>`` tree untouched. No deck uses section lists or custom shows
(verified), so there are no dangling slide-id references to clean up.

Usage:
    python scrape_shoutouts.py [--dataset dataset] [--out output/scraped]
"""
from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

# Layout names as they appear in the Google Slides export.
SECTION_HEADER_LAYOUT = "SECTION_HEADER"
SHOUTOUT_CONTENT_LAYOUT = "ONE_COLUMN_TEXT"
SHOUTOUT_HEADER_TEXT = "shoutouts"

# "General Meeting #22 4.15.pptx" -> number "22", date "4.15".
DECK_NAME_RE = re.compile(r"General Meeting #(?P<num>\d+) (?P<date>\S+)\.pptx$", re.I)


@dataclass(frozen=True)
class ScrapeResult:
    """What happened for one source deck; returned so callers/tests can assert on it."""

    source: Path
    output: Path | None  # None when the deck has no shout-out block
    kept_slide_numbers: tuple[int, ...]  # 1-based slide numbers kept from the source


# --------------------------------------------------------------------------- #
# Public interface
# --------------------------------------------------------------------------- #
def output_name_for(source_name: str) -> str:
    """Map "General Meeting #22 4.15.pptx" -> "GM #22 4.15---shoutouts.pptx".

    Raises ValueError for names that don't follow the GM convention so a stray
    file in the dataset folder fails loudly instead of producing a junk name.
    """
    m = DECK_NAME_RE.search(source_name)
    if not m:
        raise ValueError(f"Deck name does not match 'General Meeting #N date.pptx': {source_name!r}")
    return f"GM #{m['num']} {m['date']}---shoutouts.pptx"


def find_shoutout_slide_numbers(prs: PresentationType) -> list[int]:
    """Return the 1-based slide numbers of the shout-out *content* slide(s).

    Finds the "Shoutouts" section header, then collects every consecutive
    ONE_COLUMN_TEXT slide after it (GM #4 spills onto a second slide). Returns
    [] when the deck has no shout-out block (GM #1 was the first meeting).
    """
    slides = list(prs.slides)
    for idx, slide in enumerate(slides):
        if not _is_shoutout_header(slide):
            continue
        kept: list[int] = []
        for follower_idx in range(idx + 1, len(slides)):
            if slides[follower_idx].slide_layout.name != SHOUTOUT_CONTENT_LAYOUT:
                break
            kept.append(follower_idx + 1)
        return kept
    return []


def scrape_deck(source: Path, out_dir: Path) -> ScrapeResult:
    """Write the shout-out slide(s) of one deck to ``out_dir`` and report what was kept."""
    prs = Presentation(str(source))
    keep = find_shoutout_slide_numbers(prs)
    if not keep:
        return ScrapeResult(source=source, output=None, kept_slide_numbers=())
    _delete_all_slides_except(prs, keep)
    output = out_dir / output_name_for(source.name)
    prs.save(str(output))
    return ScrapeResult(source=source, output=output, kept_slide_numbers=tuple(keep))


def scrape_dataset(dataset_dir: Path, out_dir: Path) -> list[ScrapeResult]:
    """Scrape every deck in ``dataset_dir`` (numeric GM order), skipping Office lock files."""
    out_dir.mkdir(parents=True, exist_ok=True)
    decks = sorted(
        (p for p in dataset_dir.glob("*.pptx") if not p.name.startswith("~$")),
        key=_deck_sort_key,
    )
    return [scrape_deck(deck, out_dir) for deck in decks]


# --------------------------------------------------------------------------- #
# Private mechanics
# --------------------------------------------------------------------------- #
def _is_shoutout_header(slide) -> bool:
    """True for the SECTION_HEADER slide whose entire text is just "Shoutouts".

    Matching on the whole text (not a substring) matters: several decks have an
    "Updates" slide that *mentions* shout-outs, and GM #11 lists "Shoutouts for
    each week's GM!" as a to-do item.
    """
    if slide.slide_layout.name != SECTION_HEADER_LAYOUT:
        return False
    text = " ".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )
    return text.strip().lower() == SHOUTOUT_HEADER_TEXT


def _delete_all_slides_except(prs: PresentationType, keep_slide_numbers: list[int]) -> None:
    """Remove every slide not in ``keep_slide_numbers`` (1-based) from the deck.

    Two steps per slide, in this order: drop the ``<p:sldId>`` entry first, then
    drop the relationship. python-pptx's ``drop_rel`` refuses to remove a rel
    that is still referenced from the part's XML, so the sldId must go first.
    """
    keep = set(keep_slide_numbers)
    sld_id_lst = prs.slides._sldIdLst  # the only handle python-pptx exposes for this
    for number, sld_id in enumerate(list(sld_id_lst), start=1):
        if number in keep:
            continue
        r_id = sld_id.rId
        sld_id_lst.remove(sld_id)
        prs.part.drop_rel(r_id)


def _deck_sort_key(path: Path) -> int:
    """Sort decks by GM number so the run log reads #1, #2, ... not #1, #10, #11."""
    m = DECK_NAME_RE.search(path.name)
    return int(m["num"]) if m else sys.maxsize


def _main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    parser.add_argument("--dataset", type=Path, default=Path("dataset"))
    parser.add_argument("--out", type=Path, default=Path("output/scraped"))
    args = parser.parse_args(argv)

    results = scrape_dataset(args.dataset, args.out)
    for r in results:
        if r.output is None:
            print(f"SKIP  {r.source.name}: no 'Shoutouts' section found")
        else:
            print(f"OK    {r.source.name} -> {r.output.name} (source slide(s) {list(r.kept_slide_numbers)})")
    written = sum(1 for r in results if r.output)
    print(f"\n{written}/{len(results)} decks had shout-out slides; written to {args.out}")
    return 0


if __name__ == "__main__":
    sys.exit(_main())
