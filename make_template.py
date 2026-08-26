"""Build ``template/shoutouts_template.pptx`` from a scraped GM deck.

Why a template at all: generated slides should drop into the club's GM deck
looking native -- same theme, same ``ONE_COLUMN_TEXT`` layout, same ``dk1``
text colour. The cheapest way to get all of that right is to reuse their own
exported theme rather than re-create it.

What the template contains: the master and *only* the ONE_COLUMN_TEXT layout,
and zero slides. The generator adds slides from that layout. Stripping the
other ~15 layouts keeps the file small and stops PowerPoint offering a menu of
irrelevant layouts.

Usage:
    python make_template.py [--source "output/scraped/GM #22 4.15---shoutouts.pptx"]
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation

TEMPLATE_PATH = Path(__file__).resolve().parent / "template" / "shoutouts_template.pptx"
LAYOUT_NAME = "ONE_COLUMN_TEXT"
DEFAULT_SOURCE = Path("output/scraped/GM #22 4.15---shoutouts.pptx")


def build_template(source: Path, destination: Path = TEMPLATE_PATH) -> Path:
    """Strip ``source`` down to master + ONE_COLUMN_TEXT layout and save it as the template."""
    prs = Presentation(str(source))
    if len(prs.slides) == 0:
        raise SystemExit(f"{source} has no slides to take the theme from")
    used_master = prs.slides[0].slide_layout.slide_master

    # Drop every slide (sldId first, then the rel -- drop_rel refuses while referenced).
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        sld_id_lst.remove(sld_id)
        prs.part.drop_rel(r_id)

    # Google exports carry two masters (one per theme variant); keep only the one
    # the shout-out slide actually used, otherwise the layout name is ambiguous.
    master_id_lst = prs.part._element.sldMasterIdLst
    for master_id in list(master_id_lst):
        if prs.part.related_part(master_id.rId) is not used_master.part:
            master_id_lst.remove(master_id)
            prs.part.drop_rel(master_id.rId)

    # Drop every layout except the one we render on.
    kept = 0
    layout_id_lst = used_master._element.sldLayoutIdLst
    for layout_id in list(layout_id_lst):
        layout = used_master.part.related_part(layout_id.rId)
        if layout.slide_layout.name == LAYOUT_NAME:
            kept += 1
            continue
        layout_id_lst.remove(layout_id)
        used_master.part.drop_rel(layout_id.rId)
    if kept != 1:
        raise SystemExit(f"{source}: expected exactly one {LAYOUT_NAME} layout on the used master, found {kept}")

    destination.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destination))
    return destination


def _main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--out", type=Path, default=TEMPLATE_PATH)
    args = parser.parse_args(argv)
    out = build_template(args.source, args.out)
    print(f"template written: {out} ({out.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(_main())
