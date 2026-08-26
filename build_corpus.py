"""Build a spreadsheet of every shout-out text found in the scraped slides.

Why: the real shout-out form isn't available yet, so the historical slides are
the representative sample of what shout-outs look like (length, line breaks,
emoji, ALL-CAPS energy). The generator's text-measurement and layout code will
be tuned and tested against this corpus. Timestamps are deliberately omitted --
the slides never had them.

Input:  output/scraped/GM #<n> <date>---shoutouts.pptx  (from scrape_shoutouts.py)
Output: output/shoutouts_corpus.xlsx and output/shoutouts_corpus.csv
        columns: gm_number, gm_date, slide, text

Usage:
    python build_corpus.py [--scraped output/scraped] [--out output]
"""
from __future__ import annotations

import argparse
import csv
import re
import sys
from dataclasses import dataclass, fields
from pathlib import Path
from typing import Iterator

from openpyxl import Workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRAPED_NAME_RE = re.compile(r"GM #(?P<num>\d+) (?P<date>\S+)---shoutouts\.pptx$")

# python-pptx renders a soft line break (<a:br/>) as a vertical tab.
SOFT_LINE_BREAK = "\x0b"


@dataclass(frozen=True)
class Shoutout:
    """One shout-out textbox. ``slide`` is 1-based within the scraped deck."""

    gm_number: int
    gm_date: str
    slide: int
    text: str


# --------------------------------------------------------------------------- #
# Public interface
# --------------------------------------------------------------------------- #
def collect_shoutouts(scraped_dir: Path) -> list[Shoutout]:
    """Every non-empty text shape across all scraped decks, in GM / slide / z-order."""
    decks = sorted(
        (p for p in scraped_dir.glob("*---shoutouts.pptx") if SCRAPED_NAME_RE.search(p.name)),
        key=lambda p: int(SCRAPED_NAME_RE.search(p.name)["num"]),
    )
    result: list[Shoutout] = []
    for deck in decks:
        m = SCRAPED_NAME_RE.search(deck.name)
        for slide_no, slide in enumerate(Presentation(str(deck)).slides, start=1):
            for text in _iter_texts(slide.shapes):
                result.append(Shoutout(int(m["num"]), m["date"], slide_no, text))
    return result


def write_corpus(shoutouts: list[Shoutout], xlsx_path: Path, csv_path: Path) -> None:
    """Write the same rows to XLSX (for humans) and CSV (for code / diffing)."""
    header = [f.name for f in fields(Shoutout)]
    rows = [[getattr(s, name) for name in header] for s in shoutouts]

    wb = Workbook()
    ws = wb.active
    ws.title = "shoutouts"
    ws.append(header)
    for row in rows:
        ws.append(row)
    ws.column_dimensions["D"].width = 80  # text column; the rest are narrow ints
    wb.save(xlsx_path)

    # utf-8-sig so Excel opens the emoji correctly when double-clicked.
    with csv_path.open("w", newline="", encoding="utf-8-sig") as fh:
        writer = csv.writer(fh)
        writer.writerow(header)
        writer.writerows(rows)


# --------------------------------------------------------------------------- #
# Private mechanics
# --------------------------------------------------------------------------- #
def _iter_texts(shapes) -> Iterator[str]:
    """Yield cleaned text of each text-bearing shape, descending into groups.

    Soft line breaks are normalised to newlines so a cell holds ordinary
    multi-line text. Whitespace-only boxes (a few exist in the corpus) are
    dropped because they aren't shout-outs.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_texts(shape.shapes)
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.replace(SOFT_LINE_BREAK, "\n").strip()
        if text:
            yield text


def _main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    parser.add_argument("--scraped", type=Path, default=Path("output/scraped"))
    parser.add_argument("--out", type=Path, default=Path("output"))
    args = parser.parse_args(argv)

    shoutouts = collect_shoutouts(args.scraped)
    if not shoutouts:
        print(f"No scraped decks found in {args.scraped}; run scrape_shoutouts.py first", file=sys.stderr)
        return 1
    xlsx, csv_ = args.out / "shoutouts_corpus.xlsx", args.out / "shoutouts_corpus.csv"
    write_corpus(shoutouts, xlsx, csv_)
    decks = len({s.gm_number for s in shoutouts})
    print(f"{len(shoutouts)} shout-outs from {decks} decks -> {xlsx} and {csv_}")
    return 0


if __name__ == "__main__":
    sys.exit(_main())
