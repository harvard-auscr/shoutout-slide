"""End-to-end: real corpus weeks through the whole pipeline (and the CLI)."""
from __future__ import annotations

import csv
import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from shoutout_gen.deck import TEMPLATE_PATH  # noqa: E402
from shoutout_gen.generate import generate, output_name_for  # noqa: E402
from shoutout_gen.layout import overlapping_pairs  # noqa: E402

ROOT = Path(__file__).resolve().parents[1]
CORPUS = ROOT / "output" / "shoutouts_corpus.csv"
P_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

pytestmark = pytest.mark.skipif(
    not (CORPUS.exists() and TEMPLATE_PATH.exists()), reason="needs build_corpus.py and make_template.py outputs"
)


def _week(gm: int) -> list[str]:
    with CORPUS.open(encoding="utf-8-sig", newline="") as fh:
        return [r["text"] for r in csv.DictReader(fh) if int(r["gm_number"]) == gm]


@pytest.mark.parametrize("gm", [3, 22, 24])  # busiest, light, typical
def test_corpus_week_fits_one_slide_without_overlap(gm, tmp_path):
    texts = _week(gm)
    result = generate(texts, tmp_path / "w.pptx")
    assert result.slide_count == 1, result.per_slide_counts
    assert [p.text for p in result.placed] == texts  # click order preserved

    prs = Presentation(str(result.output))
    slide = prs.slides[0]
    rects = [(sh.left, sh.top, sh.width, sh.height) for sh in slide.shapes]
    assert overlapping_pairs(rects) == []
    assert all(0 <= x and x + w <= prs.slide_width and 0 <= y and y + h <= prs.slide_height for x, y, w, h in rects)
    targets = [int(t.get("spid")) for t in slide._element.findall(".//p:timing//p:spTgt", P_NS)]
    assert len(targets) == len(texts)
    assert [sh.text_frame.text for sh in slide.shapes] == texts


def test_every_corpus_week_fits_one_slide_with_auto_font_size(tmp_path):
    with CORPUS.open(encoding="utf-8-sig", newline="") as fh:
        rows = list(csv.DictReader(fh))
    sizes = {}
    for gm in sorted({int(r["gm_number"]) for r in rows}):
        texts = [r["text"] for r in rows if int(r["gm_number"]) == gm]
        result = generate(texts, tmp_path / f"{gm}.pptx")
        assert result.slide_count == 1, f"GM {gm} overflowed"
        sizes[gm] = result.font_size_pt
    assert set(sizes.values()) <= {12.0, 11.0}  # nothing in the corpus needs 10pt
    assert sum(1 for s in sizes.values() if s == 12.0) >= 20


def test_fixed_font_size_is_honoured_even_if_it_overflows(tmp_path):
    result = generate(_week(2), tmp_path / "f.pptx", font_size_pt=12)
    assert result.font_size_pt == 12 and result.slide_count == 2
    assert Presentation(str(result.output)).slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.size.pt == 12
    assert sum(result.per_slide_counts) == 42


def test_forced_compact_may_overflow_but_stays_valid(tmp_path):
    result = generate(_week(3), tmp_path / "s.pptx", font_size_pt=12, layout_mode="compact")
    assert result.slide_count >= 1 and sum(result.per_slide_counts) == 43


def test_rejects_bad_inputs(tmp_path):
    with pytest.raises(ValueError):
        generate([], tmp_path / "x.pptx")
    with pytest.raises(ValueError):
        generate(["a"], tmp_path / "x.pptx", layout_mode="diagonal")


def test_cli_end_to_end(tmp_path):
    sheet = tmp_path / "form.csv"
    with sheet.open("w", encoding="utf-8-sig", newline="") as fh:
        w = csv.writer(fh)
        w.writerow(["Timestamp", "Shout-out!"])
        for i, text in enumerate(_week(22)):
            w.writerow([f"9/10/2025 19:{i:02d}:00", text])
    proc = subprocess.run(
        [sys.executable, str(ROOT / "generate_shoutouts.py"), str(sheet), "--gm", "25", "--date", "9.10", "--out", str(tmp_path)],
        capture_output=True, text=True, cwd=ROOT,
    )
    assert proc.returncode == 0, proc.stderr
    out = tmp_path / output_name_for(25, "9.10")
    assert out.exists() and "15 shout-outs" in proc.stdout and "font 12pt" in proc.stdout
    assert [sh.text_frame.text for sh in Presentation(str(out)).slides[0].shapes] == _week(22)
    bad = subprocess.run(
        [sys.executable, str(ROOT / "generate_shoutouts.py"), str(tmp_path / "missing.csv"), "--gm", "1", "--date", "x"],
        capture_output=True, text=True, cwd=ROOT,
    )
    assert bad.returncode == 2 and "error" in bad.stderr
