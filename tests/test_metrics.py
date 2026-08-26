"""Text measurement must never under-predict what the renderer draws.

The strongest evidence available is the corpus: 641 boxes whose heights were
set by Google Slides' own autofit with real Roboto. For every box with a single
font size we re-wrap its text at its real width and demand our line count is
never lower than the rendered one (that would let boxes overlap) and usually
equal (otherwise we waste space)."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from shoutout_gen.metrics import (  # noqa: E402
    CORPUS_INSET_EMU,
    EMU_PER_PT,
    LINE_HEIGHT_FACTOR,
    MAX_BOX_WIDTH_EMU,
    TextMeasurer,
)

SCRAPED = Path(__file__).resolve().parents[1] / "output" / "scraped"

# Per-line height Google Slides used, by font size, read off the corpus.
CORPUS_LINE_HEIGHT_EMU = {10.0: 155_850, 11.0: 171_150, 12.0: 186_450}


def _corpus_boxes():
    """(font size, box width, box height, text) for single-size, non-bold corpus boxes."""
    out = []
    for deck in SCRAPED.glob("*---shoutouts.pptx"):
        for slide in Presentation(str(deck)).slides:
            for sh in slide.shapes:
                if not sh.has_text_frame or not sh.text_frame.text.strip():
                    continue
                runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
                sizes = {r.font.size.pt if r.font.size else None for r in runs}
                bold = {bool(r.font.bold) for r in runs}
                if len(sizes) == 1 and None not in sizes and bold == {False} and sizes <= set(CORPUS_LINE_HEIGHT_EMU):
                    out.append((next(iter(sizes)), sh.width, sh.height, sh.text_frame.text.replace("\x0b", "\n")))
    return out


@pytest.mark.skipif(not SCRAPED.exists(), reason="run scrape_shoutouts.py first")
def test_line_count_never_below_google_slides_rendering():
    boxes = _corpus_boxes()
    assert len(boxes) > 400
    under, exact = 0, 0
    measurers = {}
    for size, width, height, text in boxes:
        m = measurers.setdefault(size, TextMeasurer(size, inset_emu=CORPUS_INSET_EMU))
        actual_lines = round((height - 2 * CORPUS_INSET_EMU) / CORPUS_LINE_HEIGHT_EMU[size])
        predicted_lines = len(m.wrap(text, width - 2 * CORPUS_INSET_EMU))
        under += predicted_lines < actual_lines
        exact += predicted_lines == actual_lines
    assert under == 0, f"{under} boxes predicted shorter than rendered"
    assert exact / len(boxes) >= 0.85  # measured 0.885 at WIDTH_SLACK 1.10


def test_line_height_factor_is_at_least_the_corpus_value():
    for size, emu in CORPUS_LINE_HEIGHT_EMU.items():
        assert LINE_HEIGHT_FACTOR * size * EMU_PER_PT >= emu
        assert LINE_HEIGHT_FACTOR * size * EMU_PER_PT <= emu * 1.02  # but not wasteful


def test_wrap_respects_explicit_newlines_and_word_boundaries():
    m = TextMeasurer(12)
    assert m.wrap("a\nb", 10_000_000) == ["a", "b"]
    lines = m.wrap("shoutout to the best board ever", m.text_width_emu("shoutout to the") + 1)
    assert lines[0] == "shoutout to the"
    assert " ".join(lines) == "shoutout to the best board ever"


def test_overlong_word_is_split_by_character():
    m = TextMeasurer(12)
    width = m.text_width_emu("SEOUL")
    lines = m.wrap("SEOULSEOULSEOUL", width)
    assert len(lines) == 3 and "".join(lines) == "SEOULSEOULSEOUL"


def test_emoji_and_cjk_get_nonzero_fallback_width():
    m = TextMeasurer(12)
    assert m.text_width_emu("🎉") > 0
    assert m.text_width_emu("黄歆洋") > m.text_width_emu("🎉")
    # ZWJ / variation selectors add nothing on their own.
    assert m.text_width_emu("👍🏽") == m.text_width_emu("👍")


def test_measure_caps_width_and_sizes_height_by_lines():
    m = TextMeasurer(12)
    short = m.measure("bao")
    assert short.lines == ("bao",)
    assert short.height_emu == m.line_height_emu and short.width_emu < MAX_BOX_WIDTH_EMU
    long = m.measure("a very long shout-out that certainly needs to wrap onto several lines at twelve point")
    assert len(long.lines) >= 2
    assert long.height_emu == len(long.lines) * m.line_height_emu
    # Box is as wide as its widest line, never wider than the corpus cap.
    assert max(m.text_width_emu(line) for line in long.lines) == long.width_emu <= MAX_BOX_WIDTH_EMU
