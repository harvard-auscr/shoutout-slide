"""Tests for build_corpus.py: the spreadsheet must contain exactly the text boxes
on the scraped slides -- nothing dropped, nothing invented, nothing blank."""
from __future__ import annotations

import csv
import sys
from pathlib import Path

import pytest
from openpyxl import load_workbook
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
import build_corpus as bc  # noqa: E402

SCRAPED = Path(__file__).resolve().parents[1] / "output" / "scraped"
pytestmark = pytest.mark.skipif(not SCRAPED.exists(), reason="run scrape_shoutouts.py first")


@pytest.fixture(scope="module")
def shoutouts():
    return bc.collect_shoutouts(SCRAPED)


def test_every_nonblank_textbox_is_collected_once(shoutouts):
    expected = 0
    for deck in SCRAPED.glob("*---shoutouts.pptx"):
        for slide in Presentation(str(deck)).slides:
            expected += sum(1 for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip())
    assert len(shoutouts) == expected
    assert all(s.text == s.text.strip() and s.text for s in shoutouts)
    assert not any(bc.SOFT_LINE_BREAK in s.text for s in shoutouts)


def test_ordered_by_gm_number_then_slide(shoutouts):
    keys = [(s.gm_number, s.slide) for s in shoutouts]
    assert keys == sorted(keys)
    assert {s.gm_number for s in shoutouts} == set(range(2, 25))  # GM #1 has no shout-outs


def test_known_shoutouts_present(shoutouts):
    texts = {s.text for s in shoutouts}
    assert "Growl growl" in texts  # GM #22
    assert "cheese" in texts  # GM #3
    gm2_slide2 = [s for s in shoutouts if s.gm_number == 2 and s.slide == 2]
    assert len(gm2_slide2) == 7  # the overflow slide's 7 text boxes (its photo is skipped)


def test_multiline_shoutout_keeps_newlines(shoutouts):
    bees = [s for s in shoutouts if s.text.startswith("To my little bees,")]
    assert bees and "\n" in bees[0].text


def test_write_corpus_roundtrips_xlsx_and_csv(shoutouts, tmp_path):
    xlsx, csv_ = tmp_path / "c.xlsx", tmp_path / "c.csv"
    bc.write_corpus(shoutouts, xlsx, csv_)

    ws = load_workbook(xlsx)["shoutouts"]
    rows = list(ws.iter_rows(values_only=True))
    assert rows[0] == ("gm_number", "gm_date", "slide", "text")
    assert [r[3] for r in rows[1:]] == [s.text for s in shoutouts]

    with csv_.open(encoding="utf-8-sig", newline="") as fh:
        crows = list(csv.DictReader(fh))
    assert [r["text"] for r in crows] == [s.text for s in shoutouts]
    assert crows[0]["gm_number"] == "2"
