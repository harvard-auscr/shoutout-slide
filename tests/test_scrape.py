"""Tests for scrape_shoutouts.py.

Two layers:
  * pure unit tests for the naming rule (fast, no files), and
  * an integration sweep over the real ``dataset/`` decks, because the detection
    rule was derived from those decks and must keep matching every one of them
    (GM #1 has no block; GM #2/#4/#8 spill onto a second slide).
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

import pytest
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
import scrape_shoutouts as ss  # noqa: E402

ROOT = Path(__file__).resolve().parents[1]
DATASET = ROOT / "dataset"
P_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

# The decks are private and not in git, but dataset/.gitkeep IS -- so the guard
# must check for actual decks, not for the directory (which exists everywhere).
requires_dataset = pytest.mark.skipif(
    not any(DATASET.glob("General Meeting #*.pptx")), reason="private GM decks not present in dataset/"
)


# --------------------------------------------------------------------------- #
# Naming
# --------------------------------------------------------------------------- #
@pytest.mark.parametrize(
    "src, expected",
    [
        ("General Meeting #22 4.15.pptx", "GM #22 4.15---shoutouts.pptx"),
        ("General Meeting #1 9.17.pptx", "GM #1 9.17---shoutouts.pptx"),
        ("General Meeting #10 11.19.pptx", "GM #10 11.19---shoutouts.pptx"),
    ],
)
def test_output_name_for(src, expected):
    assert ss.output_name_for(src) == expected


@pytest.mark.parametrize("bad", ["random.pptx", "General Meeting 22 4.15.pptx", "GM #22.pptx"])
def test_output_name_rejects_non_gm_names(bad):
    with pytest.raises(ValueError):
        ss.output_name_for(bad)


# --------------------------------------------------------------------------- #
# Detection against the real corpus
# --------------------------------------------------------------------------- #
def _decks() -> list[Path]:
    return sorted(p for p in DATASET.glob("*.pptx") if not p.name.startswith("~$"))


# Hand-verified expectations: which source slide numbers hold the shout-outs.
EXPECTED_SLIDES = {
    1: [],  # first meeting of the year: no shout-outs yet
    2: [3, 4],  # overflow slide with 7 more shout-outs + a photo
    3: [5], 4: [4, 5], 5: [4], 6: [4], 7: [4],
    8: [4, 5],  # second slide is a photo-only shout-out
    9: [4], 10: [5], 11: [6], 12: [5], 13: [4], 14: [7], 15: [6], 16: [7],
    17: [7], 18: [7], 19: [9], 20: [8], 21: [4], 22: [4], 23: [4], 24: [4],
}


@requires_dataset
@pytest.mark.parametrize("deck", _decks(), ids=lambda p: p.name)
def test_find_shoutout_slides_matches_hand_verified_map(deck):
    num = int(re.search(r"#(\d+)", deck.name)[1])
    assert ss.find_shoutout_slide_numbers(Presentation(str(deck))) == EXPECTED_SLIDES[num]


@requires_dataset
def test_header_match_is_exact_not_substring():
    """GM #11 has a to-do slide saying 'Shoutouts for each week's GM!' -- must not match."""
    prs = Presentation(str(DATASET / "General Meeting #11 12.3.pptx"))
    assert ss.find_shoutout_slide_numbers(prs) == [6]


# --------------------------------------------------------------------------- #
# End-to-end: scrape into a temp dir and check the outputs are faithful
# --------------------------------------------------------------------------- #
@pytest.fixture(scope="module")
def scraped(tmp_path_factory):
    out = tmp_path_factory.mktemp("scraped")
    return ss.scrape_dataset(DATASET, out)


@requires_dataset
def test_scrape_dataset_writes_one_file_per_deck_with_block(scraped):
    assert len(scraped) == 24
    skipped = [r for r in scraped if r.output is None]
    assert [r.source.name for r in skipped] == ["General Meeting #1 9.17.pptx"]
    for r in scraped:
        if r.output:
            assert r.output.exists()
            assert re.fullmatch(r"GM #\d+ \S+---shoutouts\.pptx", r.output.name)


@requires_dataset
def test_scraped_slides_preserve_text_layout_and_animation_targets(scraped):
    for r in scraped:
        if r.output is None:
            continue
        src = Presentation(str(r.source))
        out = Presentation(str(r.output))
        assert len(out.slides) == len(r.kept_slide_numbers)
        assert (out.slide_width, out.slide_height) == (src.slide_width, src.slide_height)
        for out_slide, src_num in zip(out.slides, r.kept_slide_numbers):
            src_slide = src.slides[src_num - 1]
            texts = lambda s: [sh.text_frame.text for sh in s.shapes if sh.has_text_frame]
            assert texts(out_slide) == texts(src_slide)
            assert out_slide.slide_layout.name == "ONE_COLUMN_TEXT"
            # Animation survived and every target still points at a shape on the slide.
            ids = {sh.shape_id for sh in out_slide.shapes}
            targets = [int(t.get("spid")) for t in out_slide._element.findall(".//p:timing//p:spTgt", P_NS)]
            src_targets = src_slide._element.findall(".//p:timing//p:spTgt", P_NS)
            assert len(targets) == len(src_targets)
            assert all(t in ids for t in targets)
