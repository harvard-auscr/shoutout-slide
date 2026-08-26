"""Deck writer: the XML we emit must match the corpus shape and animate in click order."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from shoutout_gen.deck import TEMPLATE_PATH, PlacedShoutout, slide_size, write_deck  # noqa: E402

P_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

pytestmark = pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="run make_template.py first")


def test_template_is_16x9_with_no_slides():
    prs = Presentation(str(TEMPLATE_PATH))
    assert len(prs.slides) == 0
    assert slide_size() == (9_144_000, 5_143_500)
    assert [l.name for m in prs.slide_masters for l in m.slide_layouts] == ["ONE_COLUMN_TEXT"]


@pytest.fixture
def deck(tmp_path):
    items = [
        PlacedShoutout("first click", 0, 100_000, 100_000, 1_000_000, 200_000),
        PlacedShoutout("second click\nsecond paragraph", 0, 2_000_000, 100_000, 1_500_000, 400_000),
        PlacedShoutout("on slide two", 1, 300_000, 300_000, 900_000, 200_000),
        PlacedShoutout("third click", 0, 100_000, 900_000, 1_000_000, 200_000),
    ]
    out = write_deck(items, tmp_path / "d.pptx", font_size_pt=12)
    return items, Presentation(str(out))


def test_slides_and_shapes_match_placements(deck):
    items, prs = deck
    assert len(prs.slides) == 2
    s0 = list(prs.slides[0].shapes)
    assert [sh.text_frame.text for sh in s0] == ["first click", "second click\nsecond paragraph", "third click"]
    assert (s0[0].left, s0[0].top, s0[0].width, s0[0].height) == (100_000, 100_000, 1_000_000, 200_000)
    assert [sh.text_frame.text for sh in prs.slides[1].shapes] == ["on slide two"]
    assert all(sh.slide_layout.name == "ONE_COLUMN_TEXT" for sh in [prs.slides[0], prs.slides[1]])


def test_textbox_formatting_matches_corpus(deck):
    _, prs = deck
    shape = prs.slides[0].shapes[1]
    tf = shape.text_frame
    assert tf.word_wrap is True and tf.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    assert (tf.margin_left, tf.margin_top) == (0, 0)
    assert len(tf.paragraphs) == 2
    run = tf.paragraphs[0].runs[0]
    assert run.font.size.pt == 12
    rPr = run._r.rPr
    assert [rPr.find(f"{{{A_NS}}}{t}").get("typeface") for t in ("latin", "ea", "cs", "sym")] == ["Roboto"] * 4
    assert rPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr").get("val") == "dk1"


def test_animation_is_one_click_appear_per_shape_in_click_order(deck):
    _, prs = deck
    slide = prs.slides[0]
    el = slide._element
    effects = el.findall(".//p:timing//p:cTn[@nodeType='clickEffect']", P_NS)
    assert len(effects) == 3
    assert {e.get("presetID") for e in effects} == {"1"} and {e.get("presetClass") for e in effects} == {"entr"}
    targets = [int(t.get("spid")) for t in el.findall(".//p:timing//p:spTgt", P_NS)]
    assert targets == [sh.shape_id for sh in slide.shapes]  # click order == list order
    assert [a.text for a in el.findall(".//p:timing//p:attrName", P_NS)] == ["style.visibility"] * 3
    ids = [c.get("id") for c in el.findall(".//p:timing//p:cTn", P_NS)]
    assert len(ids) == len(set(ids))
    assert el.find(".//p:timing//p:cTn[@nodeType='mainSeq']", P_NS) is not None
