"""Write the PPTX: one textbox per shout-out, one click-triggered Appear each.

The XML emitted here is a faithful copy of what Google Slides exported for the
club's own shout-out slides (see output/scraped), so the generated slide can
be pasted into a GM deck and behave identically:

  textbox  ``<p:sp>`` with ``txBox=1``, ``wrap="square"``, ``<a:spAutoFit/>``,
           91425-EMU insets, Roboto set on latin/ea/cs/sym, colour ``dk1``.
  timing   main sequence of ``clickEffect`` nodes, preset Appear
           (presetID=1, presetClass=entr) that sets ``style.visibility`` to
           ``visible`` on the target shape. Click order == list order.

Only this module touches python-pptx / lxml; layout and measurement stay pure.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .metrics import BOX_INSET_EMU

TEMPLATE_PATH = Path(__file__).resolve().parent.parent / "template" / "shoutouts_template.pptx"
LAYOUT_NAME = "ONE_COLUMN_TEXT"
FONT_NAME = "Roboto"

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


@dataclass(frozen=True)
class PlacedShoutout:
    """A shout-out with its final geometry. ``slide`` is 0-based; EMU throughout."""

    text: str
    slide: int
    x_emu: int
    y_emu: int
    width_emu: int
    height_emu: int


# --------------------------------------------------------------------------- #
# Public interface
# --------------------------------------------------------------------------- #
def write_deck(
    shoutouts: list[PlacedShoutout],
    out_path: Path,
    font_size_pt: float,
    template_path: Path = TEMPLATE_PATH,
) -> Path:
    """Render ``shoutouts`` (in click order) into ``out_path`` and return it.

    Slides are created on demand from the template's ONE_COLUMN_TEXT layout;
    slide N of the output holds every shout-out whose ``slide == N``.
    """
    prs = Presentation(str(template_path))
    layout = _layout(prs)
    slide_count = max((s.slide for s in shoutouts), default=-1) + 1
    slides = [_blank_slide(prs, layout) for _ in range(slide_count)]

    # Shape ids per slide, in click order, for the animation sequence.
    click_order: dict[int, list[int]] = {i: [] for i in range(slide_count)}
    for item in shoutouts:
        shape = _add_textbox(slides[item.slide], item, font_size_pt)
        click_order[item.slide].append(shape.shape_id)
    for idx, slide in enumerate(slides):
        _add_click_appear_timing(slide, click_order[idx])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def slide_size(template_path: Path = TEMPLATE_PATH) -> tuple[int, int]:
    """(width, height) in EMU of the template's page, so the packer targets the real canvas."""
    prs = Presentation(str(template_path))
    return int(prs.slide_width), int(prs.slide_height)


# --------------------------------------------------------------------------- #
# Private mechanics
# --------------------------------------------------------------------------- #
def _layout(prs):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == LAYOUT_NAME:
                return layout
    raise ValueError(f"template has no {LAYOUT_NAME} layout")


def _blank_slide(prs, layout):
    """New slide from the layout with its placeholders removed (originals have none)."""
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    return slide


def _add_textbox(slide, item: PlacedShoutout, font_size_pt: float):
    shape = slide.shapes.add_textbox(
        Emu(item.x_emu), Emu(item.y_emu), Emu(item.width_emu), Emu(item.height_emu)
    )
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(BOX_INSET_EMU)

    paragraphs = item.text.split("\n")
    for i, para_text in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.space_before = paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = para_text
        run.font.size = Pt(font_size_pt)
        _set_run_font_and_colour(run)
    return shape


def _set_run_font_and_colour(run) -> None:
    """Roboto on every script slot + theme text colour, matching the corpus XML byte-for-byte."""
    rPr = run._r.get_or_add_rPr()
    fill = etree.SubElement(rPr, qn("a:solidFill"))
    etree.SubElement(fill, qn("a:schemeClr")).set("val", "dk1")
    for tag in ("a:latin", "a:ea", "a:cs", "a:sym"):
        etree.SubElement(rPr, qn(tag)).set("typeface", FONT_NAME)


def _add_click_appear_timing(slide, shape_ids_in_click_order: list[int]) -> None:
    """Append a <p:timing> main sequence with one click-Appear per shape id."""
    if not shape_ids_in_click_order:
        return
    ids = _IdCounter()
    timing = etree.SubElement(slide._element, qn("p:timing"))
    tn_lst = etree.SubElement(timing, qn("p:tnLst"))
    root_par = etree.SubElement(tn_lst, qn("p:par"))
    root = _ctn(root_par, ids, dur="indefinite", nodeType="tmRoot", restart="never")
    root_children = etree.SubElement(root, qn("p:childTnLst"))
    seq = etree.SubElement(root_children, qn("p:seq"), concurrent="1", nextAc="seek")
    main = _ctn(seq, ids, dur="indefinite", nodeType="mainSeq")
    main_children = etree.SubElement(main, qn("p:childTnLst"))
    for shape_id in shape_ids_in_click_order:
        _append_click_appear(main_children, ids, shape_id)
    for list_tag, evt in (("p:prevCondLst", "onPrev"), ("p:nextCondLst", "onNext")):
        cond = etree.SubElement(etree.SubElement(seq, qn(list_tag)), qn("p:cond"), evt=evt)
        etree.SubElement(etree.SubElement(cond, qn("p:tgtEl")), qn("p:sldTgt"))


def _append_click_appear(parent, ids: "_IdCounter", shape_id: int) -> None:
    """One click step: par(indefinite) > par(delay 0) > par(clickEffect Appear) > set visibility."""
    click = _ctn(etree.SubElement(parent, qn("p:par")), ids, fill="hold")
    _start_cond(click, "indefinite")
    group = _ctn(etree.SubElement(etree.SubElement(click, qn("p:childTnLst")), qn("p:par")), ids, fill="hold")
    _start_cond(group, "0")
    effect = _ctn(
        etree.SubElement(etree.SubElement(group, qn("p:childTnLst")), qn("p:par")),
        ids, fill="hold", nodeType="clickEffect", presetClass="entr", presetID="1", presetSubtype="0",
    )
    _start_cond(effect, "0")
    set_el = etree.SubElement(etree.SubElement(effect, qn("p:childTnLst")), qn("p:set"))
    bhvr = etree.SubElement(set_el, qn("p:cBhvr"))
    _start_cond(_ctn(bhvr, ids, dur="1", fill="hold"), "0")
    etree.SubElement(etree.SubElement(bhvr, qn("p:tgtEl")), qn("p:spTgt")).set("spid", str(shape_id))
    etree.SubElement(etree.SubElement(bhvr, qn("p:attrNameLst")), qn("p:attrName")).text = "style.visibility"
    etree.SubElement(etree.SubElement(set_el, qn("p:to")), qn("p:strVal")).set("val", "visible")


class _IdCounter:
    """Unique ``id`` for every <p:cTn>; PowerPoint expects them distinct within a slide."""

    def __init__(self) -> None:
        self.next = 1

    def take(self) -> str:
        value, self.next = self.next, self.next + 1
        return str(value)


def _ctn(parent, ids: _IdCounter, **attrs: str):
    ctn = etree.SubElement(parent, qn("p:cTn"), id=ids.take())
    for key, value in attrs.items():
        ctn.set(key, value)
    return ctn


def _start_cond(ctn, delay: str) -> None:
    etree.SubElement(etree.SubElement(ctn, qn("p:stCondLst")), qn("p:cond")).set("delay", delay)
